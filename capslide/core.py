from pptx import Presentation
import os.path
from copy import deepcopy
import json
import re

'Auto-fill subtitles into PPT'

class CapSlideError(Exception):
    """Base exception class for the CapSlide project"""
    pass

class PowerPointTemplateError(CapSlideError):
    """Base class for PowerPoint template related errors"""
    pass

class PowerPointTemplateNotFoundError(PowerPointTemplateError):
    def __init__(self, template_path):
        super().__init__(f'Template file {template_path} does not exist')
        self.template_path = template_path

    
class SubtitlesTemplateSlideIndexError(PowerPointTemplateError):
    def __init__(self, slide_index):
        if slide_index < 1:
            super().__init__(f'Slide index in template must be >= 1, current value: {slide_index}!')
        else:
            super().__init__(f'Insufficient slides in template, could not find slide at index {slide_index}!')
        self.slide_index = slide_index

class SubtitlesTemplateSlideMasterError(PowerPointTemplateError):
    def __init__(self, master_name):
        super().__init__(f'Template slide must be based on a Blank layout instead of {master_name}!')
        self.master_name = master_name

    
class SubtitlesTemplatePlaceholderError(PowerPointTemplateError):
    def __init__(self, slide_index):
        super().__init__(f'No placeholders found in slide {slide_index}, please ensure you use the #placeholder# format!')
        self.slide_index = slide_index

    
class SubtitlesProcessor:
    """
    Main functionality is to read subtitles from a file and add them to a PPT.
    Subclasses can override the unify_subtitles method to implement different subtitle processing logic.
    """
    def __init__(self, output_path, template_path, placeholder="subtitle", template_slide_page_number=0, ignore_masks=False, verbose=False) -> None:
        self.root_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        self._validate_template_pptx(template_path, template_slide_page_number)
        self.placeholder = placeholder
        self.ignore_masks = ignore_masks
        
        self.punctuation_masks = set('，。？！：；、…,.?!:;~')
        
        output_path = self.merge_path(output_path, "outputs")
        self.output_path = output_path
        self.verbose = verbose

        self.output_pptx = self.template_pptx

    def merge_path(self, path, base_dir):
        """Convert relative paths to absolute paths based on the base directory"""
        if not path.startswith('/') and not path.startswith('./'):
            path = os.path.join(self.root_dir, base_dir, path)
        return path
    
    def get_slide_by_page_number(self, pptx, slide_page_number):
        """Get slide object by slide number (1-based index)"""
        
        return pptx.slides[slide_page_number - 1]
    
    def _validate_template_pptx(self, template_path, template_slide_page_number):
        """Validate the effectiveness of the template PPTX file"""
        template_path = self.merge_path(template_path, "templates")
        
        if not os.path.exists(template_path):
            raise PowerPointTemplateNotFoundError(template_path)
        
        self.template_path = template_path
        self.template_pptx = Presentation(self.template_path)
        
        if template_slide_page_number <= 0:
            template_slide_page_number = len(self.template_pptx.slides)

        if template_slide_page_number < 0 or template_slide_page_number - 1 >= len(self.template_pptx.slides):
            raise SubtitlesTemplateSlideIndexError(template_slide_page_number)
        
        self.template_slide_page_number = template_slide_page_number
        
        template_slide = self.get_slide_by_page_number(self.template_pptx, template_slide_page_number)

        slide_layout = template_slide.slide_layout
        layout_name = slide_layout.name
        
        # Check for both English and Chinese 'Blank' layout names
        if layout_name.lower() not in ['blank', '空白']:
            raise SubtitlesTemplateSlideMasterError(layout_name)

        self.template_slide = template_slide

        matched_count = self.get_placeholders_count(self.template_slide)
        if matched_count == 0:
            raise SubtitlesTemplatePlaceholderError(template_slide_page_number)
        
        print(f'Found {matched_count} placeholders in slide index {self.template_slide_page_number} of template {self.template_path}.')

    def get_placeholders_count(self, slide, placeholder=None):
        """Count the number of placeholders in the template slide"""
        master_count = 0

        if placeholder is None:
            pattern = r'#\w+?#'
        else:
            pattern = f'#{placeholder}#'

        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        matches = re.findall(pattern, run.text)
                        master_count += len(matches)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        text_frame = cell.text_frame
                        for paragraph in text_frame.paragraphs:
                            for run in paragraph.runs:
                                matches = re.findall(pattern, run.text)
                                master_count += len(matches)
        return master_count
    
    def duplicate_slide(self, source_slide):
        """Duplicate a specified slide"""
        # Typically a blank layout
        blank_slide_layout = self.template_pptx.slide_layouts.get_by_name('Blank')
        if blank_slide_layout is None:
            blank_slide_layout = self.template_pptx.slide_layouts.get_by_name('空白')

        if blank_slide_layout is None:            
            raise Exception("Could not find a blank slide layout")

        new_slide = self.output_pptx.slides.add_slide(blank_slide_layout)

        for shape in source_slide.shapes:
            el = shape.element
            new_el = deepcopy(el)
            new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

        return new_slide

    def remove_slide(self, slide_number):
        """Remove a slide at the specified index"""
        xml_slides = self.output_pptx.slides._sldIdLst  
        slides = list(xml_slides)
        xml_slides.remove(slides[slide_number - 1])

    def replace_placeholder(self, obj, placeholder, text): 
        """Replace text in a text frame or table cell"""
        matched_count = 0
        
        text_frame = obj.text_frame
        pattern = f'#{placeholder}#'
        for p in text_frame.paragraphs:
            for r in p.runs:
                # Modifying text within the run preserves existing formatting
                if r.text.find(pattern) >= 0:
                    r.text = r.text.replace(pattern, text)
                    matched_count += 1
                    # print('Replaced text:', r.text)
        return matched_count

    def replace_placeholder_of_slide(self, slide, placeholder, text):
        """Replace text placeholders throughout the slide"""
        matched_count = 0
        if self.ignore_masks:
            text = ''.join(filter(lambda x: x not in self.punctuation_masks, text))

        for shape in slide.shapes:
            # If subtitles are in a text box, replace the text frame content
            if shape.has_text_frame:
                matched_count += self.replace_placeholder(shape, placeholder, text)
            # If subtitles are in a table, iterate through each cell
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        matched_count += self.replace_placeholder(cell, placeholder, text)
        
        return matched_count

    def append_slide_with_row(self, row):
        """Add a new slide at the end and fill it with subtitle data"""
        matched_count = 0
        new_slide = self.duplicate_slide(self.template_slide)
        for placeholder, text in row.items():
            matched_count += self.replace_placeholder_of_slide(new_slide, placeholder, text)
    
        return matched_count    

    def append_slides_with_rows(self, rows):
        """Add multiple slides at the end and fill with subtitle data"""
        total_matched_count = 0
        page_count = 0
        for row in rows:
            matched_count = self.append_slide_with_row(row)
            if matched_count > 0:
                page_count += 1
                total_matched_count += matched_count
                print(f'Successfully added page {page_count}! {matched_count} subtitle(s) inserted.')
        return total_matched_count, page_count  
    
    def append_slides_from_json_file(self, json_file_path):
        """Read subtitles from a JSON file and add them to slides"""
        json_file_path = self.merge_path(json_file_path, "data")

        if not os.path.exists(json_file_path):
            raise Exception(f'Subtitle file {json_file_path} does not exist!')
        
        with open(json_file_path, 'r') as f:
            rows = json.load(f)
            rows_count = len(rows)

            matched_count, slides_count = self.append_slides_with_rows(rows)
            
            print('*' * 40)
            print(f'Added {slides_count} of {rows_count} slides successfully! Total subtitles added: {matched_count}.')

            return matched_count, slides_count

    def append_slides_from_text_file(self, text_file_path, placeholder='subtitle'):
        """Read subtitles from a text file and add them to slides"""
        text_file_path = self.merge_path(text_file_path, "data")
        
        if not os.path.exists(text_file_path):
            raise Exception(f'Subtitle file {text_file_path} does not exist!')
        
        with open(text_file_path, 'r') as f:
            content = f.readlines()
            lines = [line.strip() for line in content if line.strip()] 
            rows = []
            for line in lines:
                row = {self.placeholder: line}
                rows.append(row)
        
            rows_count = len(rows)

            matched_count, slides_count = self.append_slides_with_rows(rows)
            
            print('*' * 40)
            print(f'Added {slides_count} of {rows_count} slides successfully! Total subtitles added: {matched_count}.')

            return matched_count, slides_count

    def append_slides_from_text_file(self, text_file_path):
        """Read subtitles from a text file and add them to slides"""
        text_file_path = self.merge_path(text_file_path, "data")
        
        if not os.path.exists(text_file_path):
            raise Exception(f'Subtitle file {text_file_path} does not exist!')
        
        with open(text_file_path, 'r') as f:
            content = f.readlines()
            lines = [line.strip() for line in content if line.strip()] 
            rows = []
            for line in lines:
                row = {self.placeholder: line}
                rows.append(row)
        
            rows_count = len(rows)

            matched_count, slides_count = self.append_slides_with_rows(rows)
            
            print('*' * 40)
            print(f'Added {slides_count} of {rows_count} slides successfully! Total subtitles added: {matched_count}.')

            return matched_count, slides_count

    def append_slides_from_file(self, file_path):
        file_path = self.merge_path(file_path, "data")
        if file_path.endswith(".json"):
            return self.append_slides_from_json_file(file_path)
        elif file_path.endswith(".txt"):
            return self.append_slides_from_text_file(file_path)
        else:
            raise Exception(f"Unsupported subtitle file type: {file_path}")

    def save(self):
        # Remove the original template slide before saving
        self.remove_slide(self.template_slide_page_number)
        self.output_pptx.save(self.output_path)
        print()
        print(f'Processing complete! File saved to {self.output_path}')